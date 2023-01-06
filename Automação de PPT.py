# -*- coding: utf-8 -*-
"""
Created on Fri Jan  6 19:46:36 2023

@author: Vinicius
"""
#importando bibliotecas
from pptx import Presentation
import os
from pptx.util import Pt
import pandas as pd

#Pegando os dados a partir da planilha
info = pd.read_excel(r'caminho\PRA FAZER ETIQUETA.xlsx')

#Adicionando texto a partir do modelo para Frase
apr = Presentation(r'caminho\Apresentação1.pptx')
a=0
b=0
while a<2:
    slide = apr.slides[a]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        
        run = p.add_run()
        run.text = info['Frase'].iloc[b]
        
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(12)
        
        b+=1
    a+=1

apr.save(r'caminho\ppt_output.pptx')
os.startfile(r'caminho\ppt_output.pptx')
