#!/usr/bin/env python
# coding: utf-8

# In[6]:


#importando bibliotecas
from pptx import Presentation
from pptx.util import Inches
import win32com.client
import os
from pptx.util import Pt
import pandas as pd


# In[3]:


#Pegando os dados a partir da planilha
info = pd.read_excel(r'C:\Users\120320\Desktop\Lumen\Impressão\PRA FAZER ETIQUETA.xlsx')


# In[4]:


len(info)


# In[101]:


info['Frase'].tail()


# In[47]:


#Adicionando texto a partir do modelo para etiqueta
apr = Presentation(r'C:\Users\120320\Desktop\Lumen\Impressão\Apresentação1.pptx')
a=0
b=0
while a<55:
    slide = apr.slides[a]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = info['Completo'].iloc[b]
        b+=1
    a+=1


# In[48]:


apr.save(r'C:\Users\120320\Desktop\Lumen\Impressão\ppt_output.pptx')
os.startfile(r'C:\Users\120320\Desktop\Lumen\Impressão\ppt_output.pptx')


# In[8]:





# In[18]:


#Adicionando texto a partir do modelo para Frase
apr = Presentation(r'C:\Users\120320\Desktop\Lumen\Impressão\Apresentação1.pptx')
a=0
b=0
while a<2:
    slide = apr.slides[a]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        
        run = p.add_run()
        run.text = info['Frase'].iloc[b]
        
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(12)
        
        b+=1
    a+=1

apr.save(r'C:\Users\120320\Desktop\Lumen\Impressão\ppt_output.pptx')
os.startfile(r'C:\Users\120320\Desktop\Lumen\Impressão\ppt_output.pptx')


# In[ ]:




